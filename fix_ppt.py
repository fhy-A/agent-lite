"""Fix last 2 slides — accurate roadmap and simple end page."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

SRC = Path.home() / "Desktop" / "Agent-Lite-产品介绍.pptx"
BLUE  = RGBColor(0x25, 0x63, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xA0, 0xBE, 0xF5)

prs = Presentation(str(SRC))

def add_text(slide, x, y, w, h, text, fs=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(fs); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = "Microsoft YaHei"; p.alignment = align
    return bx.text_frame

def add_rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

# ── Slide 15: Roadmap ──
sl = prs.slides[14]
# Clear text shapes
for s in list(sl.shapes):
    if s.has_text_frame and s.text_frame.text.strip():
        s._element.getparent().remove(s._element)

add_text(sl, 1.2, 0.5, 10, 0.6, "路线图", fs=28, bold=True)
add_rect(sl, 1.2, 1.0, 2, Pt(1.5), RGBColor(0x50, 0x85, 0xF0))

cols = [
    ("已完成 v0.5.2", "独立 Web 版", "命令安全策略\n子 Agent 并行调度\n会话分支功能\nSkill 系统 14 技能\n办公文档全支持\ni18n 中英双语\n系统托盘 + 自动更新\n412 项自动化测试"),
    ("进行中", "待办计划", "API 中转站部署\napp.js 模块化拆分"),
    ("规划中", "长期方向", "macOS / Linux 跨平台\n插件 / 扩展体系\n协作功能探索"),
]
for i, (label, title, body) in enumerate(cols):
    x = 1.2 + i * 3.7
    add_rect(sl, x, 1.8, 3.3, 4.2, RGBColor(0x2A, 0x35, 0x4A))
    add_text(sl, x+0.3, 2.0, 2.7, 0.3, label, fs=12, color=BLUE, bold=True)
    add_text(sl, x+0.3, 2.35, 2.7, 0.35, title, fs=17, bold=True)
    add_text(sl, x+0.3, 2.9, 2.7, 2.8, body, fs=13, color=RGBColor(0xB0, 0xBE, 0xD0))

# ── Slide 16: End ──
sl = prs.slides[15]
for s in list(sl.shapes):
    if s.has_text_frame and s.text_frame.text.strip():
        s._element.getparent().remove(s._element)

add_text(sl, 2, 2.2, 9, 0.9, "开始你的第一个项目", fs=44, bold=True)
add_text(sl, 2, 3.5, 9, 1.2, "下载 EXE → 双击运行 → 浏览器打开 127.0.0.1:3010\n\n用自然语言驱动代码", fs=18, color=MUTED, align=PP_ALIGN.CENTER)
add_rect(sl, 5.2, 5.5, 2.5, Pt(1.5), RGBColor(0x50, 0x85, 0xF0))
add_text(sl, 3, 5.9, 7, 0.4, "github.com/fhy-A/agent-lite", fs=13, color=MUTED, align=PP_ALIGN.CENTER)

prs.save(str(SRC))
print("Done: roadmap + end page updated")
