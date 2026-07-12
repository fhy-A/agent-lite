from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色：简约柔和风格
DARK_GRAY = RGBColor(0x2C, 0x2C, 0x2C)
SOFT_BLUE = RGBColor(0x6C, 0x8E, 0xB4)
CORAL_PINK = RGBColor(0xF2, 0xA3, 0xB6)
TEXT_GRAY = RGBColor(0x3A, 0x3A, 0x3A)
LIGHT_GRAY = RGBColor(0xF7, 0xF9, 0xFC)
BG_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
BORDER_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
SUCCESS_GREEN = RGBColor(0x5C, 0xB8, 0x5C)
WARNING_RED = RGBColor(0xE8, 0x83, 0x3A)

def add_shadow(shape):
    shadow = shape.shadow
    shadow.inherit = False
    shadow.visible = True
    shadow.distance = Pt(3)
    shadow.blur_radius = Pt(6)
    shadow.angle = 90
    shadow.alpha = 0.3

def add_shape_with_shadow(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = BORDER_GRAY
    shape.line.width = Pt(0.5)
    add_shadow(shape)
    return shape

def set_text(p, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    p.text = text
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.line_spacing = 1.4

def add_cover(title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_WHITE
    
    line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SOFT_BLUE
    line.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.333), Inches(1))
    set_text(txBox.text_frame.paragraphs[0], title, 56, DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)
    
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3.7), Inches(9.333), Inches(0.8))
    txBox2.text_frame.word_wrap = True
    set_text(txBox2.text_frame.paragraphs[0], subtitle, 20, TEXT_GRAY, align=PP_ALIGN.CENTER)
    
    txBox3 = slide.shapes.add_textbox(Inches(2), Inches(6.2), Inches(9.333), Inches(0.5))
    set_text(txBox3.text_frame.paragraphs[0], footer, 18, SOFT_BLUE, bold=True, align=PP_ALIGN.CENTER)

def add_section(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY
    
    line = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = SOFT_BLUE
    line.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9.333), Inches(1))
    set_text(txBox.text_frame.paragraphs[0], title, 44, DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)

def add_header(slide, title):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6))
    set_text(txBox.text_frame.paragraphs[0], title, 32, DARK_GRAY, bold=True)
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = SOFT_BLUE
    line.line.fill.background()

def add_simple_slide(title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_WHITE
    add_header(slide, title)
    
    y = 1.7
    for item in items:
        if item['type'] == 'heading':
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.7), Inches(0.5))
            set_text(txBox.text_frame.paragraphs[0], item['text'], 20, SOFT_BLUE, bold=True)
            y += 0.6
        elif item['type'] == 'bullet':
            txBox = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(10.7), Inches(0.4))
            set_text(txBox.text_frame.paragraphs[0], '• ' + item['text'], 16, TEXT_GRAY)
            y += 0.45
        elif item['type'] == 'text':
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.7), Inches(0.4))
            txBox.text_frame.word_wrap = True
            set_text(txBox.text_frame.paragraphs[0], item['text'], 16, TEXT_GRAY)
            y += 0.5
        elif item['type'] == 'space':
            y += 0.25

def add_comparison(title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_WHITE
    add_header(slide, title)
    
    left_card = add_shape_with_shadow(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.2), LIGHT_GRAY)
    txBox_lt = slide.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.2), Inches(0.5))
    set_text(txBox_lt.text_frame.paragraphs[0], left_title, 20, WARNING_RED, bold=True)
    
    txBox_left = slide.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(5.2), Inches(4))
    txBox_left.text_frame.word_wrap = True
    for i, item in enumerate(left_items):
        p = txBox_left.text_frame.add_paragraph() if i > 0 else txBox_left.text_frame.paragraphs[0]
        p.text = item
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_GRAY
        p.line_spacing = 1.4
        p.space_after = Pt(10)
    
    right_card = add_shape_with_shadow(slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(5.2), LIGHT_GRAY)
    txBox_rt = slide.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.2), Inches(0.5))
    set_text(txBox_rt.text_frame.paragraphs[0], right_title, 20, SUCCESS_GREEN, bold=True)
    
    txBox_right = slide.shapes.add_textbox(Inches(7.2), Inches(2.7), Inches(5.2), Inches(4))
    txBox_right.text_frame.word_wrap = True
    for i, item in enumerate(right_items):
        p = txBox_right.text_frame.add_paragraph() if i > 0 else txBox_right.text_frame.paragraphs[0]
        p.text = item
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_GRAY
        p.line_spacing = 1.4
        p.space_after = Pt(10)

def add_flow_chart(title, question, steps, value_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_WHITE
    add_header(slide, title)
    
    # 问题卡片
    q_card = add_shape_with_shadow(slide, Inches(1), Inches(2), Inches(11), Inches(0.8), RGBColor(0xFF, 0xF8, 0xE1))
    txBox = slide.shapes.add_textbox(Inches(1.3), Inches(2.2), Inches(10.4), Inches(0.4))
    set_text(txBox.text_frame.paragraphs[0], question, 18, DARK_GRAY, bold=True)
    
    # 流程步骤
    y = 3.2
    step_width = 2
    gap = 0.15
    total_width = len(steps) * step_width + (len(steps) - 1) * gap
    start_x = (13.333 - total_width) / 2
    
    for i, step in enumerate(steps):
        x = start_x + i * (step_width + gap)
        box = add_shape_with_shadow(slide, Inches(x), Inches(y), Inches(step_width), Inches(0.7), SOFT_BLUE)
        box.fill.fore_color.brightness = 0.3
        txBox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.15), Inches(step_width - 0.2), Inches(0.4))
        txBox.text_frame.word_wrap = True
        set_text(txBox.text_frame.paragraphs[0], step, 13, RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(1, Inches(x + step_width), Inches(y + 0.3), Inches(gap), Inches(0.1))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = SOFT_BLUE
            arrow.line.fill.background()
    
    # 价值说明
    v_card = add_shape_with_shadow(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(2), RGBColor(0xF0, 0xF9, 0xFF))
    txBox = slide.shapes.add_textbox(Inches(1.8), Inches(4.7), Inches(9.4), Inches(1.6))
    txBox.text_frame.word_wrap = True
    for i, line in enumerate(value_text):
        p = txBox.text_frame.add_paragraph() if i > 0 else txBox.text_frame.paragraphs[0]
        p.text = line
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_GRAY
        p.line_spacing = 1.3
        p.space_after = Pt(8)

def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_WHITE
    add_header(slide, title)
    
    # 简化表格用卡片展示关键对比
    cols = 3
    card_width = 3.6
    x_positions = [0.8, 4.85, 8.9]
    
    for i in range(min(cols, len(headers))):
        card = add_shape_with_shadow(slide, Inches(x_positions[i]), Inches(1.8), Inches(card_width), Inches(5.2), LIGHT_GRAY)
        
        txBox = slide.shapes.add_textbox(Inches(x_positions[i] + 0.2), Inches(2), Inches(card_width - 0.4), Inches(0.5))
        set_text(txBox.text_frame.paragraphs[0], headers[i], 18, SOFT_BLUE, bold=True)
        
        txBox2 = slide.shapes.add_textbox(Inches(x_positions[i] + 0.2), Inches(2.6), Inches(card_width - 0.4), Inches(4.2))
        txBox2.text_frame.word_wrap = True
        for j, row in enumerate(rows[i]):
            p = txBox2.text_frame.add_paragraph() if j > 0 else txBox2.text_frame.paragraphs[0]
            p.text = row
            p.font.name = 'Microsoft YaHei'
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_GRAY
            p.line_spacing = 1.3
            p.space_after = Pt(6)

# 第1页：封面
add_cover('Agent Lite', '本地运行的 AI 编程助手\n直接操作项目文件，数据不离开电脑', 'github.com/fhy-A/agent-lite')

# 第2页：定位
add_simple_slide('Agent Lite 是什么？', [
    {'type': 'heading', 'text': '目标用户'},
    {'type': 'bullet', 'text': '独立开发者 / 小团队 / 重视数据隐私的程序员'},
    {'type': 'space'},
    {'type': 'heading', 'text': '产品定位'},
    {'type': 'bullet', 'text': '不是 VS Code 插件 — 不绑定编辑器'},
    {'type': 'bullet', 'text': '不是 SaaS 服务 — 不传数据到云端'},
    {'type': 'bullet', 'text': '不是 CLI 工具 — 有完整图形界面'},
    {'type': 'space'},
    {'type': 'text', 'text': '一句话：浏览器打开的本地 AI 助手，能直接读写项目文件'},
    {'type': 'space'},
    {'type': 'heading', 'text': '核心理念'},
    {'type': 'bullet', 'text': '本地优先 — 100% 运行在 127.0.0.1，数据不离开电脑'},
    {'type': 'bullet', 'text': '安全可控 — 权限分级 + 操作审批 + 自动备份'},
    {'type': 'bullet', 'text': '零门槛 — 打开浏览器即用，无需安装 IDE、配置插件'}
])

# 第3页：痛点 vs 解法
add_comparison('痛点 vs 解法', '通用 AI 的 6 个痛点',
    ['只能聊天，不能操作文件 → 来回复制粘贴',
     '让它改文件 → 输出整段代码让你自己动手',
     '问了 50 轮上下文爆炸 → 开始胡言乱语',
     '多任务只能串行 → 一个等一个',
     '试了新方案想回退 → 只能新建聊天重来',
     '数据经过第三方服务器 → 公司代码不敢贴'],
    'Agent Lite 的 6 个解法',
    ['直接在项目目录里读写、搜索、执行命令',
     '生成 diff → 你审批 → 自动写入，原件备份',
     '超 95% 自动压缩为摘要，保留关键信息',
     '一次启动 3 个子 Agent 并行分析不同模块',
     '当前消息一键分叉，独立探索后随时切回',
     '100% 本地运行，数据不离开你的电脑'])

# 第4页：场景一 - 代码重构
add_flow_chart('场景一：代码重构', '📋 "帮我把 app.js 里超过 200 行的函数拆小"',
    ['search_files\n定位', 'read_file\n读取', 'propose_edit\n生成 diff', '你审批', 'apply_edit\n写入+备份', 'pytest\n验证'],
    ['💡 从提出需求到代码落盘，全程在浏览器内闭环',
     '不像 Chat：你复制一段 → 它输出修改 → 你粘贴回去 → 漏了换行还要排查'])

# 第5页：场景二 - 数据分析
add_flow_chart('场景二：数据分析', '📋 "分析 sales.xlsx 各产品线季度趋势，生成 Word 报告"',
    ['openpyxl\n读取', 'pandas\n透视', 'python-docx\n生成报告', 'output/\n保存', '文件树\n预览'],
    ['💡 Excel → 分析 → 报告，全程不离开 Agent Lite',
     '不像 Chat：上传 → 输出 Markdown 表格 → 你手动排版',
     '支持 docx / xlsx / pptx / pdf / csv 读 + 写 + 创建'])

# 第6页：场景三 - 并行探索
add_flow_chart('场景三：并行探索', '📋 "同时检查 app.js、server.py、styles.css 的潜在问题"',
    ['主 Agent\n拆任务', '子 Agent 1\n并行', '子 Agent 2\n并行', '子 Agent 3\n并行', '主 Agent\n合并'],
    ['💡 3 个文件的审查同时完成，不是串行等',
     '每个子 Agent 独立上下文 + 独立用量跟踪，完成后精确合并'])

# 第7页：场景四 - 方案对比
add_flow_chart('场景四：方案对比', '📋 "Flask 改 FastAPI 试试？不好再切回来"',
    ['Branches\n点击', '新建分支', '上下文\n自动复制', '独立探索', '不满意\n秒切回'],
    ['💡 分支独立保存，互不影响。支持多层嵌套 A→B→C',
     '不像 Chat：新建聊天 → 重新描述背景 → 还是丢了之前的上下文'])

# 第8页：Skill 系统
add_simple_slide('Skill 系统（14 个内置技能）', [
    {'type': 'text', 'text': '可复用的自动化任务模板，Agent 根据对话自动匹配并执行'},
    {'type': 'space'},
    {'type': 'bullet', 'text': 'office-files — Word/Excel/PPT/PDF 读写'},
    {'type': 'bullet', 'text': 'code-review / receiving-code-review — 代码审查全流程'},
    {'type': 'bullet', 'text': 'test-driven-development — TDD 驱动开发'},
    {'type': 'bullet', 'text': 'python-testing — Python 测试框架'},
    {'type': 'bullet', 'text': 'design-aesthetics — UI 设计规范'},
    {'type': 'bullet', 'text': 'brainstorming / writing-plans / executing-plans — 规划与执行'},
    {'type': 'bullet', 'text': '支持用户自定义 Skill（Markdown 编写，0 代码）'},
    {'type': 'space'},
    {'type': 'text', 'text': '💡 Skill = 可复用的能力模板，让 Agent 越用越懂你的项目'}
])

# 第9页：记忆系统
add_simple_slide('记忆系统', [
    {'type': 'text', 'text': '跨会话、跨项目保留关键信息'},
    {'type': 'space'},
    {'type': 'heading', 'text': '自动提取'},
    {'type': 'bullet', 'text': '对话结束后 Agent 主动提炼关键信息'},
    {'type': 'space'},
    {'type': 'heading', 'text': '手动管理'},
    {'type': 'bullet', 'text': 'Memory 面板增删改查'},
    {'type': 'space'},
    {'type': 'heading', 'text': '智能匹配'},
    {'type': 'bullet', 'text': '新对话开始时根据上下文自动注入相关记忆'},
    {'type': 'space'},
    {'type': 'text', 'text': '💡 Memory = 跨会话的知识库，两者结合让 Agent 越用越懂你的项目'}
])

# 第10页：权限体系
add_table_slide('权限体系：人机协作的安全模型',
    ['Plan 计划模式', 'Accept 审批模式（默认）', 'Bypass 自动执行'],
    [['只读 + 生成方案', '修改操作需要审批', '', '适合场景：', '• 调研分析', '• 代码审查', '• 方案探索'],
     ['读 + 写 + 执行', '文件修改弹窗确认', '', '适合场景：', '• 日常开发', '• 功能迭代', '• Bug 修复'],
     ['全自动免审批', '跳过确认直接写入', '', '适合场景：', '• 完全信任的脚本', '• 批量自动化任务', '• 紧急修复']])

# 第11页：安全保障
add_simple_slide('你的代码安全吗？', [
    {'type': 'text', 'text': '当 AI 能操作你的文件时，安全不是"信不信它"，而是"你能控到哪一步"'},
    {'type': 'space'},
    {'type': 'heading', 'text': '🛡 命令拦截'},
    {'type': 'bullet', 'text': 'rm -rf / → 拦截 | format C: → 拦截 | 删除系统文件 → 拦截'},
    {'type': 'space'},
    {'type': 'heading', 'text': '🛡 文件保护'},
    {'type': 'bullet', 'text': '项目外路径 → 自动重定向 | 每次写入前自动备份'},
    {'type': 'space'},
    {'type': 'heading', 'text': '🛡 网络隔离'},
    {'type': 'bullet', 'text': '内网 IP 拦截 | 127.0.0.1 访问拦截'},
    {'type': 'space'},
    {'type': 'heading', 'text': '🛡 提示注入防护 + 每步可审可拒可回滚'}
])

# 第12页：产品对比
add_simple_slide('Agent Lite 的差异化优势', [
    {'type': 'heading', 'text': '4 个"唯一"'},
    {'type': 'space'},
    {'type': 'bullet', 'text': '✓ 唯一无需安装任何运行时（浏览器即用）'},
    {'type': 'bullet', 'text': '✓ 唯一 100% 本地存储（代码不经第三方服务器）'},
    {'type': 'bullet', 'text': '✓ 唯一支持会话分支（试错成本最低）'},
    {'type': 'bullet', 'text': '✓ 唯一完全免费开源'},
    {'type': 'space'},
    {'type': 'space'},
    {'type': 'text', 'text': '定位：不是功能最多的，但是门槛最低、最轻量、数据最安全的选择'}
])

# 第13页：技术栈
add_table_slide('技术栈',
    ['前端', '后端', '质量保障'],
    [['零前端依赖', '无 npm / webpack', '纯 JS 12,000 行', '单文件 · 零构建'],
     ['Python stdlib', 'http.server 原生', '可选 docx / openpyxl', '无 Flask / Django'],
     ['412 测试 · 20 秒', '全量自动化', 'CI 就绪', '零外部依赖']])

# 第14页：路线图
add_simple_slide('路线图', [
    {'type': 'heading', 'text': '✅ 已完成 v0.5.2'},
    {'type': 'bullet', 'text': '独立 Web 版 · 命令安全策略 · 子 Agent 并行调度'},
    {'type': 'bullet', 'text': '会话分支功能 · Skill 系统 14 技能 · 办公文档全支持'},
    {'type': 'bullet', 'text': 'i18n 中英双语 · 412 项自动化测试'},
    {'type': 'space'},
    {'type': 'heading', 'text': '🚧 进行中'},
    {'type': 'bullet', 'text': 'API Key 自主配置 · 自由切换模型 · 深色/浅色主题 · 一键 EXE 分发'},
    {'type': 'space'},
    {'type': 'heading', 'text': '🔮 规划中'},
    {'type': 'bullet', 'text': '更多模型适配 · macOS / Linux 支持 · 插件/扩展体系 · 协作功能探索'},
    {'type': 'space'},
    {'type': 'text', 'text': '💡 支持自主配置 API Key 和 Base URL，接入任意兼容 OpenAI 协议的模型服务'}
])

# 第15页：结尾
add_cover('开始使用', 'cd 到项目目录 → 运行 agent-lite\n浏览器打开 127.0.0.1:3010\n用自然语言驱动代码', 'github.com/fhy-A/agent-lite')

prs.save('Agent-Lite-产品介绍-完整版.pptx')
print('PPT final: 15 slides with flow charts')
