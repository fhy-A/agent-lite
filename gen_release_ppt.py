# -*- coding: utf-8 -*-
"""将 release-notes.md 转换为 PPT，每个版本区块一个 slide。"""
import os
import re
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(ROOT, "release-notes.md")
OUT = os.path.join(os.path.expanduser("~"), "Desktop", "release-notes.pptx")

ACCENT = RGBColor(0x66, 0x7e, 0xea)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x55, 0x55, 0x55)


def parse_sections(text):
    """按 '## ' 标题切分为版本区块，忽略 '---' 分隔符。"""
    lines = text.splitlines()
    sections = []
    current = None
    for line in lines:
        if line.strip() == "---":
            continue
        m = re.match(r"^##\s+(.*)", line)
        if m:
            if current:
                sections.append(current)
            current = {"title": m.group(1).strip(), "body": []}
        elif current is not None:
            current["body"].append(line)
    if current:
        sections.append(current)
    return sections


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景色
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = ACCENT

    tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Agent Lite"
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Release Notes"
    r2.font.size = Pt(28)
    r2.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    return slide


def add_section_slide(prs, section):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = section["title"]
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = ACCENT

    # 内容
    body_lines = [ln for ln in section["body"] if ln.strip()]
    cb = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(9), Inches(5.5))
    ctf = cb.text_frame
    ctf.word_wrap = True
    first = True
    if not body_lines:
        p = ctf.paragraphs[0]
        r = p.add_run()
        r.text = "(无详细说明)"
        r.font.size = Pt(16)
        r.font.color.rgb = GRAY
        r.font.italic = True
        return slide

    for line in body_lines:
        stripped = line.strip()
        sub = re.match(r"^###\s+(.*)", stripped)
        h4 = re.match(r"^####\s+(.*)", stripped)
        if first:
            p = ctf.paragraphs[0]
            first = False
        else:
            p = ctf.add_paragraph()

        if sub or h4:
            heading = (sub or h4).group(1).strip()
            r = p.add_run()
            r.text = heading
            r.font.size = Pt(18)
            r.font.bold = True
            r.font.color.rgb = DARK
            p.space_before = Pt(8)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            content = clean_md(stripped[2:])
            p.level = 1
            r = p.add_run()
            r.text = "• " + content
            r.font.size = Pt(14)
            r.font.color.rgb = GRAY
        else:
            r = p.add_run()
            r.text = clean_md(stripped)
            r.font.size = Pt(14)
            r.font.color.rgb = GRAY
    return slide


def clean_md(s):
    """去除简单的 markdown 标记：**bold**、`code`。"""
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"`(.+?)`", r"\1", s)
    return s


def main():
    with open(SRC, "r", encoding="utf-8") as f:
        text = f.read()
    sections = parse_sections(text)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    for sec in sections:
        add_section_slide(prs, sec)

    prs.save(OUT)
    print("生成完成：{}".format(OUT))
    print("版本区块数：{}，总 slide 数：{}".format(len(sections), len(sections) + 1))


if __name__ == "__main__":
    main()
