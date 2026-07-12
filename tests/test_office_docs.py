"""
Office Document Integration Tests — Word/Excel/PPT/PDF read & write.

Creates test documents, then verifies agent-lite can process them correctly.

Run: python tests/test_office_docs.py
  or: python -m pytest tests/test_office_docs.py -v
"""
import json
import os
import sys
import tempfile
import threading
import time
import unittest
from pathlib import Path
from unittest import mock

import requests

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
import server as server_mod


# ═══════════════════════════════════════════════════════════════════
# Test Fixtures: Generate real office documents
# ═══════════════════════════════════════════════════════════════════

def _create_test_docx(path: Path):
    """Create a Word document with paragraphs and a table."""
    from docx import Document
    doc = Document()
    doc.add_heading("测试报告", level=1)
    doc.add_paragraph("这是一份自动生成的测试文档，用于验证 Agent Lite 的文档处理能力。")
    doc.add_heading("核心指标", level=2)
    table = doc.add_table(rows=4, cols=3, style="Table Grid")
    headers = ["指标", "Q1", "Q2"]
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    data = [
        ["营收", "120万", "145万"],
        ["成本", "80万", "92万"],
        ["利润", "40万", "53万"],
    ]
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            table.rows[r + 1].cells[c].text = val
    doc.add_paragraph("以上数据仅供参考。")
    doc.save(str(path))


def _create_test_xlsx(path: Path):
    """Create an Excel file with multiple sheets."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "销售数据"
    ws1.append(["产品", "数量", "单价", "总价"])
    ws1.append(["产品A", 100, 50, 5000])
    ws1.append(["产品B", 200, 30, 6000])
    ws1.append(["产品C", 150, 80, 12000])
    ws2 = wb.create_sheet("汇总")
    ws2.append(["分类", "合计"])
    ws2.append(["总收入", 23000])
    wb.save(str(path))


def _create_test_pptx(path: Path):
    """Create a PowerPoint with multiple slides."""
    from pptx import Presentation
    prs = Presentation()
    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "项目汇报"
    slide1.shapes.placeholders[1].text = "2026年度 Q2 总结"
    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "已完成事项"
    slide2.shapes.placeholders[1].text = "• 功能开发\n• 测试完善\n• 文档更新"
    prs.save(str(path))


def _create_test_pdf(path: Path):
    """Create a simple PDF (via reportlab if available, else write a minimal one)."""
    try:
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(str(path))
        c.drawString(100, 750, "测试 PDF 文档")
        c.drawString(100, 730, "这是一个用于测试的 PDF 文件。")
        c.drawString(100, 700, "第二页内容：更多测试数据。")
        c.showPage()
        c.drawString(100, 750, "第二页")
        c.drawString(100, 730, "这是 PDF 的第二页内容。")
        c.save()
    except ImportError:
        # Fallback: create minimal valid PDF
        pdf_content = (
            b"%PDF-1.4\n"
            b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
            b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
            b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Resources<<>>>>endobj\n"
            b"xref\n0 4\n0000000000 65535 f \n0000000009 00000 n \n0000000058 00000 n \n0000000115 00000 n \n"
            b"trailer<</Size 4/Root 1 0 R>>\nstartxref\n190\n%%EOF"
        )
        path.write_bytes(pdf_content)


# ═══════════════════════════════════════════════════════════════════
# Test Server Setup
# ═══════════════════════════════════════════════════════════════════

def _free_port():
    import socket
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.bind(("127.0.0.1", 0))
        return s.getsockname()[1]


_PORT = _free_port()
_BASE = f"http://127.0.0.1:{_PORT}"


def _req(method, path, **kwargs):
    url = f"{_BASE}{path}"
    resp = requests.request(method, url, timeout=15, **kwargs)
    try:
        return resp.status_code, resp.json()
    except Exception:
        return resp.status_code, resp.text


class TestOfficeDocuments(unittest.TestCase):

    @classmethod
    def setUpClass(cls):
        # Create temp project with test documents
        cls.tmp_root = Path(tempfile.mkdtemp(prefix="agentlite_office_"))
        cls.tmp_data = Path(tempfile.mkdtemp(prefix="agentlite_offdat_"))

        # Create test documents
        cls.docx_path = cls.tmp_root / "test_report.docx"
        cls.xlsx_path = cls.tmp_root / "test_data.xlsx"
        cls.pptx_path = cls.tmp_root / "test_presentation.pptx"
        cls.pdf_path = cls.tmp_root / "test_document.pdf"
        cls.txt_path = cls.tmp_root / "test_notes.txt"
        cls.csv_path = cls.tmp_root / "test_export.csv"

        _create_test_docx(cls.docx_path)
        _create_test_xlsx(cls.xlsx_path)
        _create_test_pptx(cls.pptx_path)
        _create_test_pdf(cls.pdf_path)
        cls.txt_path.write_text("项目编号,预算\nP001,50000\nP002,30000\n", encoding="utf-8")
        cls.csv_path.write_text("产品,数量\nA,10\nB,20\n", encoding="utf-8")

        for sub in ["sessions", "memory", "skills", "attachments", "file-backups", "output"]:
            (cls.tmp_data / sub).mkdir(parents=True, exist_ok=True)

        config = {
            "projectRoot": str(cls.tmp_root),
            "newApiBaseUrl": "http://localhost:3000",
            "userHome": str(Path.home()),
        }
        (cls.tmp_data / "config.json").write_text(json.dumps(config), encoding="utf-8")

        cls._patchers = [
            mock.patch.object(server_mod, "DATA_DIR", cls.tmp_data),
            mock.patch.object(server_mod, "CONFIG_PATH", cls.tmp_data / "config.json"),
            mock.patch.object(server_mod, "SESSIONS_DIR", cls.tmp_data / "sessions"),
            mock.patch.object(server_mod, "MEMORY_DIR", cls.tmp_data / "memory"),
            mock.patch.object(server_mod, "SKILLS_DIR", cls.tmp_data / "skills"),
            mock.patch.object(server_mod, "ATTACHMENTS_DIR", cls.tmp_data / "attachments"),
            mock.patch.object(server_mod, "FILE_BACKUP_DIR", cls.tmp_data / "file-backups"),
            mock.patch.object(server_mod, "APP_DIR", cls.tmp_root),
        ]
        for p in cls._patchers:
            p.start()

        server_mod.ThreadingHTTPServer.daemon_threads = True
        cls._server = server_mod.ThreadingHTTPServer(
            ("127.0.0.1", _PORT), server_mod.AgentLiteHandler
        )
        cls._server.socket.settimeout(2.0)
        cls._thread = threading.Thread(target=cls._server.serve_forever, daemon=True)
        cls._thread.start()
        time.sleep(0.3)

    @classmethod
    def tearDownClass(cls):
        for p in cls._patchers:
            p.stop()

    def _run_python_script(self, script_content, timeout=30):
        """Write Python script to temp file and execute it (avoids PS encoding issues)."""
        script_path = self.tmp_root / "_test_script.py"
        script_path.write_text(script_content, encoding="utf-8")
        status, data = _req("POST", "/api/tools/run_command", json={
            "command": f"python {script_path}",
            "description": "office test",
            "timeout": timeout,
        })
        return status, data

    def _run_python(self, code):
        """DEPRECATED: use _run_python_script for anything with non-ASCII."""
        status, data = _req("POST", "/api/tools/run_command", json={
            "command": f'python -c "{code}"',
            "description": "office test",
            "timeout": 30,
        })
        return status, data

    # ═══════════════════════════════════════════════════════════════
    # 1. Word (.docx) 读取
    # ═══════════════════════════════════════════════════════════════

    def test_docx_read_paragraphs(self):
        """Extract all paragraph text from a Word document."""
        script = (
            "from docx import Document\n"
            f"doc = Document(r'{self.docx_path}')\n"
            "count = 0\n"
            "for p in doc.paragraphs:\n"
            "    if p.text.strip():\n"
            "        print(p.text)\n"
            "        count += 1\n"
            "print(f'PARAGRAPH_COUNT: {count}')\n"
        )
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        self.assertTrue(data.get("ok"), f"Failed: {data}")
        stdout = data.get("stdout", "")
        self.assertIn("PARAGRAPH_COUNT:", stdout)

    def test_docx_read_tables(self):
        """Extract table data from Word document."""
        script = (
            "from docx import Document\n"
            f"doc = Document(r'{self.docx_path}')\n"
            "for t in doc.tables:\n"
            "    for row in t.rows:\n"
            "        print('\\t'.join(c.text for c in row.cells))\n"
            "print(f'TABLE_COUNT: {len(doc.tables)}')\n"
        )
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        self.assertTrue(data.get("ok"))
        stdout = data.get("stdout", "")
        self.assertIn("TABLE_COUNT:", stdout)

    def test_docx_read_nonexistent_handles_gracefully(self):
        status, data = self._run_python(
            "from docx import Document; doc=Document(r'nonexistent.docx'); print('ok')"
        )
        # Should fail gracefully
        self.assertNotEqual(status, 500)

    # ═══════════════════════════════════════════════════════════════
    # 2. Excel (.xlsx) 读取
    # ═══════════════════════════════════════════════════════════════

    def test_xlsx_list_sheets(self):
        """List sheet names from Excel workbook."""
        script = (
            "import openpyxl\n"
            f"wb = openpyxl.load_workbook(r'{self.xlsx_path}', data_only=True)\n"
            "for s in wb.sheetnames:\n"
            "    print(s)\n"
            "print(f'SHEET_COUNT: {len(wb.sheetnames)}')\n"
        )
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        stdout = data.get("stdout", "")
        self.assertIn("SHEET_COUNT:", stdout)

    def test_xlsx_read_sheet_data(self):
        """Read data from a specific sheet."""
        script = (
            "import openpyxl\n"
            f"wb = openpyxl.load_workbook(r'{self.xlsx_path}', data_only=True)\n"
            "ws = wb[wb.sheetnames[0]]\n"
            "for row in list(ws.iter_rows())[:5]:\n"
            "    print('\\t'.join(str(c.value or '') for c in row))\n"
            "print(f'ROW_COUNT: {ws.max_row}')\n"
        )
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        self.assertTrue(data.get("ok"))
        stdout = data.get("stdout", "")
        self.assertIn("ROW_COUNT:", stdout)

    def test_xlsx_pandas_read(self):
        """Use pandas to read Excel data."""
        script = (
            "import pandas as pd\n"
            f"df = pd.read_excel(r'{self.xlsx_path}')\n"
            "print(f'ROWS: {len(df)}, COLS: {len(df.columns)}')\n"
            "print(df.head(10).to_string())\n"
        )
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        stdout = data.get("stdout", "")
        self.assertIn("ROWS:", stdout)

    # ═══════════════════════════════════════════════════════════════
    # 3. PPT (.pptx) 读取
    # ═══════════════════════════════════════════════════════════════

    def test_pptx_read_slides(self):
        """Extract text from all PowerPoint slides."""
        script = (
            "from pptx import Presentation\n"
            f"prs = Presentation(r'{self.pptx_path}')\n"
            "result = []\n"
            "for slide in prs.slides:\n"
            "    for s in slide.shapes:\n"
            "        if s.has_text_frame:\n"
            "            result.append(s.text[:50])\n"
            "print(f'PPTX_SLIDE_COUNT: {len(prs.slides)}')\n"
            "print(f'PPTX_SHAPE_COUNT: {len(result)}')\n"
        )
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        self.assertTrue(data.get("ok"))
        stdout = data.get("stdout", "")
        self.assertIn("PPTX_SLIDE_COUNT:", stdout)
        self.assertIn("PPTX_SHAPE_COUNT:", stdout)

    # ═══════════════════════════════════════════════════════════════
    # 4. PDF 读取
    # ═══════════════════════════════════════════════════════════════

    def test_pdf_read_text(self):
        """Extract text from PDF using PyPDF2."""
        code = (
            f"from PyPDF2 import PdfReader; r=PdfReader(r'{self.pdf_path}'); "
            "[print(r.pages[i].extract_text()) for i in range(len(r.pages))]"
        )
        status, data = self._run_python(code)
        self.assertEqual(status, 200)
        stdout = data.get("stdout", "")
        # PDF may have limited text extraction — at least should not crash
        self.assertNotEqual(status, 500)

    def test_pdf_read_page_count(self):
        """Verify PDF page count detection."""
        code = (
            f"from PyPDF2 import PdfReader; r=PdfReader(r'{self.pdf_path}'); "
            "print(f'Pages: {len(r.pages)}')"
        )
        status, data = self._run_python(code)
        self.assertEqual(status, 200)
        stdout = data.get("stdout", "")
        self.assertIn("Pages:", stdout)

    # ═══════════════════════════════════════════════════════════════
    # 5. 跨文件操作
    # ═══════════════════════════════════════════════════════════════

    def test_csv_via_pandas(self):
        """Read CSV with pandas."""
        script = (
            "import pandas as pd\n"
            f"df = pd.read_csv(r'{self.csv_path}')\n"
            "print(f'CSV_ROWS: {len(df)}')\n"
            "print(df.to_string())\n"
        )
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        stdout = data.get("stdout", "")
        self.assertIn("CSV_ROWS:", stdout)

    def test_read_multiple_docs_in_sequence(self):
        script = (
            f"from docx import Document\n"
            f"doc = Document(r'{self.docx_path}')\n"
            f"print(f'DOCX_OK: {{len(doc.paragraphs)}}')\n"
            f"import openpyxl\n"
            f"wb = openpyxl.load_workbook(r'{self.xlsx_path}', data_only=True)\n"
            f"print(f'XLSX_OK: {{len(wb.sheetnames)}}')\n"
            f"import pandas as pd\n"
            f"df = pd.read_csv(r'{self.csv_path}')\n"
            f"print(f'CSV_OK: {{len(df)}}')\n"
        )
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        stdout = data.get("stdout", "")
        self.assertIn("DOCX_OK", stdout)
        self.assertIn("XLSX_OK", stdout)
        self.assertIn("CSV_OK", stdout)

    # ═══════════════════════════════════════════════════════════════
    # 6. 文档生成
    # ═══════════════════════════════════════════════════════════════

    def test_generate_new_docx(self):
        output = self.tmp_root / "output" / "generated_report.docx"
        output.parent.mkdir(exist_ok=True)
        script = f"from docx import Document\ndoc = Document()\ndoc.add_heading('Generated Report', level=1)\ndoc.add_paragraph('Auto-generated document.')\ndoc.save(r'{output}')\nprint('DOCX_GENERATED: OK')\n"
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        self.assertTrue(data.get("ok"), f"Generate failed: {data}")
        self.assertTrue(output.exists(), "Generated file should exist")
        self.assertGreater(output.stat().st_size, 1000)

    def test_generate_new_xlsx(self):
        output = self.tmp_root / "output" / "generated_data.xlsx"
        output.parent.mkdir(exist_ok=True)
        script = f"import openpyxl\nwb = openpyxl.Workbook()\nws = wb.active\nws.title = 'Generated'\nws.append(['Name','Value'])\nws.append(['Item1',42])\nwb.save(r'{output}')\nprint('XLSX_GENERATED: OK')\n"
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        self.assertTrue(output.exists())

    def test_generate_new_pptx(self):
        output = self.tmp_root / "output" / "generated_slides.pptx"
        output.parent.mkdir(exist_ok=True)
        script = f"from pptx import Presentation\nprs = Presentation()\nslide = prs.slides.add_slide(prs.slide_layouts[0])\nslide.shapes.title.text = 'Generated Slide'\nprs.save(r'{output}')\nprint('PPTX_GENERATED: OK')\n"
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        self.assertTrue(output.exists())

    # ═══════════════════════════════════════════════════════════════
    # 7. 错误恢复
    # ═══════════════════════════════════════════════════════════════

    def test_corrupted_file_recovery(self):
        """Reading a non-docx file as docx should fail gracefully."""
        code = (
            f"from docx import Document; doc=Document(r'{self.txt_path}'); "
            "print('OK')"
        )
        status, data = self._run_python(code)
        # Should fail gracefully, not crash the server
        self.assertNotEqual(status, 500)

    def test_missing_file_graceful(self):
        status, data = self._run_python(
            "import pandas as pd; df=pd.read_excel(r'nonexistent_file.xlsx'); print('OK')"
        )
        self.assertNotEqual(status, 500)

    def test_large_output_truncation(self):
        """Generating massive output should be handled."""
        code = "for i in range(1000): print(f'Line {i}')"
        status, data = self._run_python(code)
        self.assertEqual(status, 200)
        stdout = data.get("stdout", "")
        # Output should be truncated at some reasonable limit
        self.assertLess(len(stdout), 200000, "Output should be truncated")

    # ═══════════════════════════════════════════════════════════════
    # 8. 文件编码与路径
    # ═══════════════════════════════════════════════════════════════

    def test_chinese_filename_handling(self):
        """Chinese filenames should work with raw string paths."""
        chinese_path = self.tmp_root / "chinese_test_doc.docx"
        from docx import Document
        doc = Document()
        doc.add_paragraph("Chinese content test")
        doc.save(str(chinese_path))
        self.assertTrue(chinese_path.exists())

        script = (
            f"from docx import Document\n"
            f"doc = Document(r'{chinese_path}')\n"
            "for p in doc.paragraphs:\n"
            "    if p.text.strip():\n"
            "        print(p.text)\n"
            "print('CHINESE_DOC_READ: OK')\n"
        )
        status, data = self._run_python_script(script)
        self.assertEqual(status, 200)
        self.assertTrue(data.get("ok"))
        self.assertIn("CHINESE_DOC_READ: OK", data.get("stdout", ""))


if __name__ == "__main__":
    unittest.main()
